import os
import re
from pptx import Presentation
from docx import Document

def clean_xml_string(s):
    # Supprime les caractères de contrôle ASCII (0–8, 11–12, 14–31) et NULL (\x00)
    return re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F]', '', s)

def normalize_line_endings(s):
    # Remplace tous les types de sauts de ligne par un seul \n
    # Attention, il reste certains sauts de ligne dans PowerPoint qui sont affichés collés dans la string renvoyée, pour une raison que j'ignore
    return re.sub(r'\r\n?|\u0085|\u2028|\u2029', '\n', s)

def list_pptx_files():
    """Lists all .pptx files in the current directory."""
    return [f for f in os.listdir() if f.endswith(".pptx")]

def choose_pptx_file():
    """Prompts the user to select a .pptx file from the current directory."""
    pptx_files = list_pptx_files()
    
    if not pptx_files:
        print("No PowerPoint files found in the current directory.")
        return None

    print("Select a PowerPoint file:")
    for i, file in enumerate(pptx_files, 1):
        print(f"{i}. {file}")

    while True:
        try:
            choice = int(input("Enter the number of the file: "))
            if 1 <= choice <= len(pptx_files):
                return pptx_files[choice - 1]
            else:
                print("Invalid choice. Please enter a valid number.")
        except ValueError:
            print("Invalid input. Please enter a number.")

def extract_script_from_notes(pptx_path, output_filename):
    """Extracts script content from speaker notes based on custom rules."""
    prs = Presentation(pptx_path)
    doc = Document()
    table = doc.add_table(rows=1, cols=4)
    table.style = 'Table Grid'
    header_cells = table.rows[0].cells
    header_cells[0].text = "Code VO"
    header_cells[1].text = "Script"
    header_cells[2].text = "Voix"
    header_cells[3].text = "Prononciation"

    for i, slide in enumerate(prs.slides, start=1):
        if not slide.notes_slide or not slide.notes_slide.notes_text_frame.text.strip():
            continue

        notes_text = clean_xml_string(normalize_line_endings(slide.notes_slide.notes_text_frame.text)).strip()
        print(repr(slide.notes_slide.notes_text_frame.text))
        matches = re.findall(r"## (.*?)\n(.*?)(?=\n\n|\Z)", notes_text, re.DOTALL)

        if not matches:
            continue

        subindex = 1

        for name, lines in matches:
            for line in lines.strip().split("\n"):
                row_cells = table.add_row().cells
                row_cells[0].text = f"{i}.{subindex}"
                row_cells[1].text = line.strip()
                row_cells[2].text = name.strip()
                subindex += 1

        doc.save(output_filename)

def main():
    print("PowerPoint Script Extractor")
    print("---------------------------")
    print("Warning: This tool will overwrite existing Word files. Please close any open Word documents before proceeding.")
    pptx_file = choose_pptx_file()
    
    if pptx_file:
        output_filename = f"{os.path.splitext(pptx_file)[0]}_script.docx"
        extract_script_from_notes(pptx_file, output_filename)

        print(f"Script extracted successfully! Saved as {output_filename}")
        print("Press Enter to exit.")
        input()

if __name__ == "__main__":
    main()
